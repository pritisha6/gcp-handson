"""Document parsing, chunking, and embedding pipeline.

Extracts text from uploaded documents (PDF, PPTX, XLSX, HTML, TXT, CSV),
splits it into overlapping token-bounded chunks, and upserts those chunks
into Pinecone (via ``PineconeClient``) for later retrieval-augmented
design generation.
"""
import os
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import tiktoken
from bs4 import BeautifulSoup
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError

from app.config import Settings, get_settings
from app.db.pinecone_client import PineconeClient, get_pinecone_client
from app.schemas.document import DocumentChunk
from app.utils.errors import DocumentProcessingError, FileTooLargeError, UnsupportedFileTypeError
from app.utils.logger import get_logger

logger = get_logger(__name__)

_SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".xlsx", ".html", ".htm", ".txt", ".csv"}
_TOKEN_ENCODING = "cl100k_base"


class DocumentProcessor:
    """Parses supported document formats into token-chunked, embeddable text."""

    def __init__(
        self,
        pinecone_client: Optional[PineconeClient] = None,
        settings: Optional[Settings] = None,
    ) -> None:
        self._settings = settings or get_settings()
        self._pinecone_client = pinecone_client or get_pinecone_client()
        self._encoding = tiktoken.get_encoding(_TOKEN_ENCODING)

    # --- Format-specific extraction ---

    def process_pdf(self, file_path: str) -> List[str]:
        """Extract text from a PDF, one entry per page.

        Args:
            file_path: Path to the PDF file on disk.

        Returns:
            A list of page texts (empty string for pages with no extractable text).

        Raises:
            DocumentProcessingError: If the file is corrupt or unreadable.
            FileTooLargeError: If the file exceeds the configured size limit.
        """
        self._check_file_size(file_path)
        try:
            reader = PdfReader(file_path)
            return [page.extract_text() or "" for page in reader.pages]
        except (PdfReadError, OSError, ValueError) as exc:
            raise DocumentProcessingError(f"Could not parse PDF '{file_path}': {exc}") from exc

    def process_pptx(self, file_path: str) -> List[str]:
        """Extract text from a PPTX, one entry per slide.

        Args:
            file_path: Path to the PPTX file on disk.

        Returns:
            A list of slide texts (all text frames on a slide, joined with newlines).

        Raises:
            DocumentProcessingError: If the file is corrupt or unreadable.
            FileTooLargeError: If the file exceeds the configured size limit.
        """
        self._check_file_size(file_path)
        try:
            presentation = Presentation(file_path)
            slides_text: List[str] = []
            for slide in presentation.slides:
                pieces = [
                    shape.text_frame.text
                    for shape in slide.shapes
                    if shape.has_text_frame and shape.text_frame.text
                ]
                slides_text.append("\n".join(pieces))
            return slides_text
        except Exception as exc:
            raise DocumentProcessingError(f"Could not parse PPTX '{file_path}': {exc}") from exc

    def process_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract sheet data from an XLSX workbook.

        Args:
            file_path: Path to the XLSX file on disk.

        Returns:
            One dict per sheet: ``{"sheet_name": str, "headers": List[str], "rows": List[Dict]}``.
            The first non-empty row is treated as the header row.

        Raises:
            DocumentProcessingError: If the file is corrupt or unreadable.
            FileTooLargeError: If the file exceeds the configured size limit.
        """
        self._check_file_size(file_path)
        try:
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            sheets: List[Dict[str, Any]] = []
            for sheet in workbook.worksheets:
                headers: List[str] = []
                data_rows: List[Dict[str, Any]] = []
                for row in sheet.iter_rows(values_only=True):
                    if row is None or all(cell is None for cell in row):
                        continue
                    if not headers:
                        headers = [str(cell) if cell is not None else f"col_{j}" for j, cell in enumerate(row)]
                        continue
                    data_rows.append(
                        {headers[j] if j < len(headers) else f"col_{j}": cell for j, cell in enumerate(row)}
                    )
                sheets.append({"sheet_name": sheet.title, "headers": headers, "rows": data_rows})
            workbook.close()
            return sheets
        except Exception as exc:
            raise DocumentProcessingError(f"Could not parse Excel file '{file_path}': {exc}") from exc

    def process_html(self, file_path: str) -> str:
        """Extract visible text from an HTML file.

        Args:
            file_path: Path to the HTML file on disk.

        Returns:
            Whitespace-normalized visible text content (script/style tags excluded).

        Raises:
            DocumentProcessingError: If the file cannot be read.
            FileTooLargeError: If the file exceeds the configured size limit.
        """
        self._check_file_size(file_path)
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            return "\n".join(line.strip() for line in text.splitlines() if line.strip())
        except OSError as exc:
            raise DocumentProcessingError(f"Could not read HTML file '{file_path}': {exc}") from exc

    def process_text_or_csv(self, file_path: str) -> str:
        """Read a plain text or CSV file as-is.

        Args:
            file_path: Path to the TXT/CSV file on disk.

        Returns:
            The file's text content.

        Raises:
            DocumentProcessingError: If the file cannot be read.
            FileTooLargeError: If the file exceeds the configured size limit.
        """
        self._check_file_size(file_path)
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except OSError as exc:
            raise DocumentProcessingError(f"Could not read file '{file_path}': {exc}") from exc

    # --- Orchestration ---

    def process_file(self, file_path: str, upload_to_pinecone: bool = True) -> Dict[str, Any]:
        """Auto-detect a file's type, extract text, chunk it, and (optionally) index it.

        Args:
            file_path: Path to the uploaded file on disk.
            upload_to_pinecone: Whether to embed and upsert chunks to Pinecone.

        Returns:
            A dict matching ``ProcessedDocument``'s shape: ``document_id``,
            ``filename``, ``source_type``, ``total_chunks``, ``chunks``, ``warnings``.

        Raises:
            UnsupportedFileTypeError: If the extension is not supported.
            DocumentProcessingError: If parsing fails.
            FileTooLargeError: If the file exceeds the configured size limit.
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        if extension not in _SUPPORTED_EXTENSIONS:
            raise UnsupportedFileTypeError(path.name, extension)

        document_id = str(uuid.uuid4())
        warnings: List[str] = []
        segments: List[Dict[str, Any]] = []

        if extension == ".pdf":
            source_type = "pdf"
            for i, page_text in enumerate(self.process_pdf(file_path)):
                if page_text.strip():
                    segments.append({"text": page_text, "metadata": {"page": i + 1}})
                else:
                    warnings.append(f"Page {i + 1} had no extractable text.")
        elif extension == ".pptx":
            source_type = "pptx"
            for i, slide_text in enumerate(self.process_pptx(file_path)):
                if slide_text.strip():
                    segments.append({"text": slide_text, "metadata": {"slide": i + 1}})
                else:
                    warnings.append(f"Slide {i + 1} had no extractable text.")
        elif extension == ".xlsx":
            source_type = "xlsx"
            for sheet in self.process_excel(file_path):
                sheet_text = self._sheet_to_text(sheet)
                if sheet_text.strip():
                    segments.append({"text": sheet_text, "metadata": {"sheet": sheet["sheet_name"]}})
                else:
                    warnings.append(f"Sheet '{sheet['sheet_name']}' had no data.")
        elif extension in (".html", ".htm"):
            source_type = "html"
            text = self.process_html(file_path)
            if text.strip():
                segments.append({"text": text, "metadata": {}})
            else:
                warnings.append("No visible text found in HTML document.")
        else:  # .txt / .csv
            source_type = "csv" if extension == ".csv" else "txt"
            text = self.process_text_or_csv(file_path)
            if text.strip():
                segments.append({"text": text, "metadata": {}})
            else:
                warnings.append("File was empty.")

        chunks = self._build_chunks(segments, document_id=document_id, filename=path.name, source_type=source_type)

        if upload_to_pinecone and chunks:
            try:
                self._pinecone_client.upsert_documents(
                    [
                        {
                            "chunk_id": chunk.chunk_id,
                            "text": chunk.text,
                            "document_id": chunk.document_id,
                            "filename": chunk.filename,
                            "source_type": chunk.source_type,
                            "chunk_index": chunk.chunk_index,
                            **chunk.metadata,
                        }
                        for chunk in chunks
                    ]
                )
            except Exception:
                logger.exception("Failed to upload chunks for document '%s' to Pinecone", document_id)
                warnings.append("Chunks were extracted but could not be uploaded to the vector store.")

        logger.info(
            "Processed document '%s' (%s): %d chunks, %d warnings",
            path.name,
            source_type,
            len(chunks),
            len(warnings),
        )

        return {
            "document_id": document_id,
            "filename": path.name,
            "source_type": source_type,
            "total_chunks": len(chunks),
            "chunks": chunks,
            "warnings": warnings,
        }

    # --- Internal helpers ---

    def _check_file_size(self, file_path: str) -> None:
        max_bytes = self._settings.MAX_UPLOAD_FILE_SIZE_MB * 1024 * 1024
        try:
            size = os.path.getsize(file_path)
        except OSError as exc:
            raise DocumentProcessingError(f"Could not access file '{file_path}': {exc}") from exc
        if size > max_bytes:
            raise FileTooLargeError(Path(file_path).name, size, max_bytes)

    def _sheet_to_text(self, sheet: Dict[str, Any]) -> str:
        headers = sheet["headers"]
        lines = [", ".join(headers)]
        for row in sheet["rows"]:
            lines.append(", ".join(str(row.get(h, "")) for h in headers))
        return "\n".join(lines)

    def _build_chunks(
        self,
        segments: List[Dict[str, Any]],
        *,
        document_id: str,
        filename: str,
        source_type: str,
    ) -> List[DocumentChunk]:
        chunk_size = self._settings.CHUNK_SIZE_TOKENS
        overlap = self._settings.CHUNK_OVERLAP_TOKENS
        chunks: List[DocumentChunk] = []
        index = 0
        for segment in segments:
            for chunk_text, token_count in self._chunk_text(segment["text"], chunk_size, overlap):
                chunks.append(
                    DocumentChunk(
                        chunk_id=f"{document_id}-{index}",
                        document_id=document_id,
                        filename=filename,
                        source_type=source_type,
                        chunk_index=index,
                        text=chunk_text,
                        token_count=token_count,
                        metadata=segment["metadata"],
                    )
                )
                index += 1
        return chunks

    def _chunk_text(self, text: str, chunk_size: int, overlap: int) -> List[Tuple[str, int]]:
        """Split text into overlapping token windows.

        Args:
            text: Source text to split.
            chunk_size: Target tokens per chunk.
            overlap: Token overlap between consecutive chunks.

        Returns:
            A list of (chunk_text, token_count) tuples.
        """
        tokens = self._encoding.encode(text)
        if not tokens:
            return []

        step = max(chunk_size - overlap, 1)
        windows: List[Tuple[str, int]] = []
        for start in range(0, len(tokens), step):
            window = tokens[start : start + chunk_size]
            if not window:
                break
            windows.append((self._encoding.decode(window), len(window)))
            if start + chunk_size >= len(tokens):
                break
        return windows


_document_processor: Optional[DocumentProcessor] = None


def get_document_processor() -> DocumentProcessor:
    """Return a process-wide singleton DocumentProcessor (FastAPI dependency)."""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor
